"""Build the Lighthouse idea-submission deck from the Redrob template, export PDF.

  python scripts/build_deck.py [--runtime "Xs"]

Fills the template's section slides with real content + the generated architecture
diagram, inserts a dedicated Problem slide, saves a .pptx, and converts to PDF via
PowerPoint COM (Windows). All numbers come from eval/results.md.
"""
import argparse
import copy
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def duplicate_slide(prs, index):
    """Deep-copy a slide (shapes + image relationships) and append it."""
    src = prs.slides[index]
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    rid_map = {}
    for rid, rel in src.part.rels.items():
        if not rel.is_external and "image" in rel.reltype:
            rid_map[rid] = new.part.relate_to(rel.target_part, rel.reltype)
    for blip in new.shapes._spTree.iter(qn("a:blip")):
        old = blip.get(qn("r:embed"))
        if old in rid_map:
            blip.set(qn("r:embed"), rid_map[old])
    return new


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)

TEMPLATE = r"C:/Users/Utkarsh/Downloads/Idea Submission Template _ Redrob.pptx"
OUT_PPTX = os.path.abspath("deck/Lighthouse_Redrob_Deck.pptx")
OUT_PDF = os.path.abspath("deck/Lighthouse_Redrob_Deck.pdf")
DIAGRAM = os.path.abspath("deck/assets/architecture.png")
INK = RGBColor(0x1F, 0x2A, 0x44)


def _no_inherited_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(parse_xml(f'<a:buNone xmlns:a="{A_NS}"/>'))


def set_body(slide, bullets, size=12, body_shape=None):
    """Replace the body (longest text box, not the title) with bullet content.
    We suppress the placeholder's inherited bullet and use a consistent '• '."""
    if body_shape is None:
        cands = [sh for sh in slide.shapes if sh.has_text_frame]
        body_shape = max(cands, key=lambda s: len(s.text_frame.text)) if cands else None
    if body_shape is None:
        return
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _no_inherited_bullet(p)
        run = p.add_run()
        run.text = ("• " if lvl == 0 else "    – ") + txt
        run.font.size = Pt(size - 1 if lvl else size)
        run.font.color.rgb = INK
        p.space_after = Pt(5)


def add_textbox(slide, left, top, width, height, text, size, bold=False, color=INK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def title_of(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and 0 < len(sh.text_frame.text) < 50:
            return sh.text_frame.text.strip()
    return ""


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--runtime", default="< 5 min")
    args = ap.parse_args()

    prs = Presentation(TEMPLATE)

    # --- create a Problem slide by duplicating a CONTENT slide (clean background) ---
    dup = duplicate_slide(prs, 1)                 # clone "Solution Overview" layout/bg
    move_slide(prs, len(prs.slides._sldIdLst) - 1, 1)   # place it right after the title

    slides = list(prs.slides)
    # the duplicate (now index 1) still carries the source title/body text; retitle it
    prob_title = min((sh for sh in slides[1].shapes if sh.has_text_frame),
                     key=lambda s: len(s.text_frame.text))
    prob_title.text_frame.paragraphs[0].runs[0].text = "The Problem"
    by_title = {title_of(s): s for s in slides[2:]}   # map the original section slides

    # --- Title slide: clean tagline + problem statement (no overlay on the hero art) ---
    title = slides[0]
    for sh in title.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("Team Name"):
            sh.text_frame.paragraphs[0].runs[0].text = "Team Name :  (your team)"
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("Problem Statement"):
            sh.text_frame.paragraphs[0].runs[0].text = (
                "Problem Statement : Rank the top 100 of 100,000 candidates for one nuanced "
                "Senior AI Engineer JD with recruiter-grade judgment — not keyword matching — "
                "avoiding keyword-stuffer traps and ~80 planted honeypots.   "
                "Project: Lighthouse — a reasoning-first, explainable candidate ranker.")

    # --- Problem slide content ---
    problem = slides[1]
    set_body(problem, [
        ("100,000 candidates; genuine Senior-AI-Engineer fits are rare — only 0.78% even have an AI-ish title.", 0),
        ("Seeded traps: keyword-stuffers (an Accountant listing 9 AI skills), plain-language Tier-5s (real systems, no buzzwords), behavioral twins, and ~80 honeypots (subtly impossible profiles).", 0),
        ("The provided sample_submission ranks the stuffers #1–20 — the exact wrong answer. >10% honeypots in the top-100 = disqualification.", 0),
        ("Scoring is NDCG@10-heavy: the top-10 must be right. The real task is reasoning about the gap between what the JD says and what it means.", 0),
    ])

    content = {
        "Solution Overview": [
            ("Lighthouse — an offline-LLM-augmented hybrid ranker: recruiter-grade, reasoning-first, fully explainable.", 0),
            ("Five fit components (semantic + role-coherence + career-evidence + experience + skill-trust), gated by JD hard-negatives and a behavioral modifier, with honeypots zeroed.", 0),
            ("Reads career evidence & semantics, not keywords — surfaces plain-language Tier-5s and sinks keyword-stuffers.", 0),
            ("Every pick ships a grounded, no-hallucination reason.", 0),
            ("CPU-only, < 5 min, no API at rank-time — production-shaped, not a GPT-per-candidate demo.", 0),
        ],
        "JD Understanding & Candidate Evaluation": [
            ("Must-haves: production embeddings/retrieval, vector/hybrid search ops, ranking-eval (NDCG/MRR/MAP), strong Python. Ideal 6–8 yrs, 4–5 applied-ML at product (not services) cos, shipped a ranking/search/recsys system.", 0),
            ("Hard-negatives encoded as gates: services-only, outside-India & won't-relocate (no visa), research-only, CV/speech-only w/o NLP, LangChain-only-recent, title-chasers, non-engineering current role.", 0),
            ("Signal priority: career-history evidence > job title > listed skills. Behavioral signals decide actual availability.", 0),
            ("Claude (offline) parsed the JD into a static committed rubric — read at rank-time with no API call.", 0),
        ],
        "Ranking Methodology": [
            ("Retrieve: precomputed bge-small embeddings + BM25 over clean per-candidate text blobs (100K).", 0),
            ("Score, each in [0,1]: semantic_fit (cosine vs 10 JD facets), role_coherence (taxonomy+semantic), career_evidence (built ranking/search at product cos), experience_fit (soft 6–8 curve), trust_skills (proficiency×duration×endorsements×assessment).", 0),
            ("Combine: weighted base (role_coherence 0.26 & career_evidence 0.24 lead) × Π(hard-negative gates) × behavioral modifier [0.80–1.12]; honeypots → 0.", 0),
            ("final = base × gates × behavior. Order by score desc; ties → candidate_id ascending; take top 100.", 0),
        ],
        "Explainability & Data Validation": [
            ("Deterministic, fact-grounded reasoning per candidate — cites real years, title, named skills, employers, signal values, and names honest gaps (e.g., '120-day notice; Toronto-based — relocation risk').", 0),
            ("Zero hallucination: every term traces to the candidate's own profile (enforced by tests). Variation by dominant factor; tone tracks the rank band.", 0),
            ("Honeypot/anomaly detector (explainable rules): ≥3 advanced/expert skills with 0 months, date contradictions, tenure overflowing experience, invalid education — flags 0.042% of the pool, no false positives on real profiles.", 0),
            ("Schema-tolerant parsing; sentinels (-1 github/offer, empty assessments — 60–76% of pool) treated as neutral, never penalized. Fully seeded/deterministic.", 0),
        ],
        "End-to-End Workflow": [
            ("Phase 1 — Offline precompute (network OK, no time limit): JD →(Claude)→ jd_rubric.json; candidates → text blobs → bge-small embeddings (fp16) + BM25; Claude → eval labels. Artifacts committed.", 0),
            ("Phase 2 — Rank-time (CPU, no network, < 5 min): load artifacts → score 5 components → gates × behavioral → honeypot zero → sort/tie-break → top-100 → grounded reasoning → submission.csv.", 0),
            ("Single reproduce command: python rank.py --candidates ./data/candidates.jsonl --out ./submission.csv → passes the official validator.", 0),
        ],
        "Results & Performance": [
            ("Label-INDEPENDENT evidence (no trust in our labels needed):", 0),
            ("0 honeypots in the top-100 (audited over full 100K; DQ threshold >10%). All 100/100 hold an AI/ML/IR/DS/Search/NLP title — 0 non-technical; the sample_submission instead ranks HR/Accountants at #1–20.", 1),
            ("Trap resistance: the keyword baseline puts 13/32 keyword-stuffers in its top-25; Lighthouse admits 0. Remove the anti-trap stack and honeypots climb (median rank 209→143).", 1),
            ("Directional only (self-labeled — NOT absolute accuracy): vs keyword baseline composite 0.998 vs 0.553. Labeler & ranker share assumptions, so a ~perfect NDCG@10 = internal consistency, not validated accuracy — the gap + ablation deltas are the signal. A blind human-label harness ships for an independent check.", 0),
            (f"Constraints met: rank step CPU-only, no network, {args.runtime} over 100K; embeddings precomputed (~75 MB fp16); 36 tests green.", 0),
        ],
        "Technologies Used": [
            ("sentence-transformers BAAI/bge-small-en-v1.5 — small, CPU-friendly, strong retrieval quality; precomputed so rank-time needs no model.", 0),
            ("rank_bm25 — lexical recall signal complementing dense retrieval.", 0),
            ("numpy / pandas — the entire rank step: fast, deterministic, dependency-light (fits the 5-min / 16 GB / CPU box).", 0),
            ("Claude (offline only) — JD→rubric parsing and eval labeling; never called at rank-time, no candidate data sent to any API.", 0),
            ("pytest (36 tests), Streamlit (HF sandbox), python-pptx + matplotlib (this deck).", 0),
        ],
        "Submission Assets": [
            ("GitHub: https://github.com/jacklachan/lighthouse-indiaruns", 0),
            ("Live sandbox (HF Spaces): https://huggingface.co/spaces/jacklachan/lighthouse", 0),
            ("Reproduce: pip install -r requirements.txt → python precompute.py → python rank.py → python validate_submission.py submission.csv", 0),
            ("Demo video: [placeholder — add link before submission]", 0),
            ("AI use: Claude offline for JD parsing + eval labeling + coding assistant; no candidate data hit any API at rank-time (see submission_metadata.yaml).", 0),
        ],
    }

    for t, bullets in content.items():
        if t in by_title:
            set_body(by_title[t], bullets)

    # --- architecture diagram ---
    arch = by_title.get("System Architecture")
    if arch and os.path.exists(DIAGRAM):
        arch.shapes.add_picture(DIAGRAM, Inches(0.35), Inches(1.5), width=Inches(9.3))

    os.makedirs("deck", exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"Saved {OUT_PPTX}")

    # --- PDF via PowerPoint COM ---
    try:
        import win32com.client
        pp = win32com.client.Dispatch("PowerPoint.Application")
        deck = pp.Presentations.Open(OUT_PPTX, WithWindow=False)
        deck.SaveAs(OUT_PDF, 32)  # 32 = ppSaveAsPDF
        deck.Close()
        pp.Quit()
        print(f"Saved {OUT_PDF}")
    except Exception as e:
        print(f"PDF export skipped ({e}). Open the .pptx and Save As PDF.")


if __name__ == "__main__":
    main()
